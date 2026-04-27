"""Generate project presentation PPT with correct slide ordering."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

WHITE = RGBColor(255,255,255)
BLACK = RGBColor(0,0,0)
DARK = RGBColor(30,30,30)
GRAY = RGBColor(100,100,100)
BLUE = RGBColor(37,99,235)
GREEN = RGBColor(22,163,74)
RED = RGBColor(220,38,38)
PURPLE = RGBColor(147,51,234)
ACCENT = RGBColor(37,99,235)
FONT_NAME = "Times New Roman"

def add_bg(slide, color=WHITE):
    bg = slide.background; fill = bg.fill; fill.solid(); fill.fore_color.rgb = color

def add_text(slide, left, top, width, height, text, size=18, bold=False, color=BLACK, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text; p.font.size = Pt(size); p.font.bold = bold
    p.font.color.rgb = color; p.font.name = FONT_NAME; p.alignment = align
    return txBox

def add_bullet_slide(slide, left, top, width, height, items, size=16, color=BLACK):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item; p.font.size = Pt(size); p.font.color.rgb = color
        p.font.name = FONT_NAME; p.space_after = Pt(8)

def add_line(slide, top):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(top), Inches(11.7), Pt(2))
    shape.fill.solid(); shape.fill.fore_color.rgb = ACCENT; shape.line.fill.background()

def add_title(slide, title_text, color=DARK):
    add_text(slide, 0.8, 0.4, 11.7, 0.7, title_text, size=32, bold=True, color=color)
    add_line(slide, 1.1)

def add_table(slide, left, top, width, height, headers, rows, font_size=12):
    """Add a proper pptx table to a slide."""
    num_rows = len(rows) + 1
    num_cols = len(headers)
    table_shape = slide.shapes.add_table(num_rows, num_cols, Inches(left), Inches(top), Inches(width), Inches(height))
    tbl = table_shape.table
    # Header row
    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.text = h
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(font_size)
            paragraph.font.bold = True
            paragraph.font.name = FONT_NAME
            paragraph.font.color.rgb = WHITE
            paragraph.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(30, 30, 30)
    # Data rows
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = tbl.cell(i + 1, j)
            cell.text = str(val)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(font_size)
                paragraph.font.name = FONT_NAME
                paragraph.alignment = PP_ALIGN.CENTER
            if i % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(240, 240, 255)
    return table_shape

graphs_dir = os.path.join(os.path.dirname(__file__), "results", "graphs")
dist_graphs_dir = os.path.join(os.path.dirname(__file__), "results", "distributed_graphs")

# ═══ SLIDE 1: Title ═══
sl = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(sl); add_line(sl, 2.8)
add_text(sl, 0.8, 1.0, 11.7, 1.5, "Performance Evaluation of Encryption\nAlgorithms on Edge Devices", size=36, bold=True, color=DARK, align=PP_ALIGN.CENTER)
add_text(sl, 0.8, 3.2, 11.7, 0.6, "AES-256-GCM  |  ChaCha20-Poly1305  |  PRESENT  |  SPECK-64/128", size=20, color=BLUE, align=PP_ALIGN.CENTER)
add_text(sl, 0.8, 4.5, 11.7, 0.5, "Distributed Systems - Final Project Presentation", size=18, color=GRAY, align=PP_ALIGN.CENTER)

# ═══ SLIDE 2: Problem Statement ═══
sl = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(sl); add_title(sl, "Problem Statement")
add_bullet_slide(sl, 0.8, 1.5, 11.7, 5.5, [
    "- Security Constraint: IoT edge devices (sensors, smart cameras, medical monitors) collect sensitive data that must be encrypted immediately before transmission.",
    "- Hardware Constraint: Edge devices possess strictly limited CPU cycles, memory, and battery life.",
    "- Conventional Problem: Traditional standard encryption (like full hardware AES) may exhaust the resources of ultra-constrained devices (e.g., simple RFID tags or microcontrollers).",
    "",
    "Core Objective:",
    "To conduct a rigorous comparative evaluation of standard versus lightweight encryption algorithms -- measuring latency, throughput, CPU utilization, and memory footprint -- to identify the optimal algorithm for diverse edge environments."
], size=18)

# ═══ SLIDE 3: Justification ═══
sl = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(sl); add_title(sl, "Justification for Problem Statement")
add_bullet_slide(sl, 0.8, 1.5, 11.7, 5.5, [
    "Why is this problem important?",
    "",
    "1. Explosive IoT Growth: By 2025, over 75 billion IoT devices are projected worldwide. Each device generates sensitive data requiring encryption.",
    "",
    "2. Edge Security Gap: Most IoT breaches occur at the edge. Encrypting data at the source (not in the cloud) prevents man-in-the-middle interception during transmission.",
    "",
    "3. Resource Mismatch: Standard AES requires dedicated hardware (AES-NI) that passive RFID tags and 8-bit microcontrollers simply do not have. Lightweight ciphers (PRESENT, SPECK) were specifically designed for these constraints.",
    "",
    "4. No Universal Solution: There is no single 'best' algorithm -- the optimal choice depends on whether the device has hardware acceleration, how much RAM/ROM is available, and the acceptable latency budget.",
    "",
    "5. Distributed Systems Relevance: In a real deployment, thousands of edge nodes must encrypt concurrently. Our benchmark measures not just individual cipher performance but system-wide scalability under parallel distributed workloads."
], size=16)

# ═══ SLIDE 4: Literature Survey (Table 1/2) ═══
sl = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(sl); add_title(sl, "Literature Survey (1/2)")
add_table(sl, 0.5, 1.3, 12.3, 5.5,
    ["#", "Paper Title", "Authors / Year", "Key Contribution / Relevance"],
    [
        ["1", "The Design of Rijndael: AES", "Daemen & Rijmen, 2002", "Establishes AES block cipher. We use AES-256-GCM as industry-standard baseline."],
        ["2", "ChaCha, a variant of Salsa20", "Bernstein D.J., 2008", "ChaCha20 stream cipher excels in software without hardware crypto instructions."],
        ["3", "PRESENT: An Ultra-Lightweight Block Cipher", "Bogdanov et al., CHES 2007", "SPN cipher for extreme constraints (~1,570 GE). Benchmarks lowest hardware tier."],
        ["4", "The SIMON and SPECK Families", "Beaulieu et al., IACR 2013", "SPECK-64/128 ARX cipher by NSA. 24 cycles/byte on ARM Cortex microcontrollers."],
        ["5", "Lightweight Crypto for IoT Devices", "Thakor et al., IEEE Access 2021", "Comparative framework for latency, throughput, memory. Our methodology mirrors this."],
    ], font_size=11)

# ═══ SLIDE 4b: Literature Survey (Table 2/2) ═══
sl = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(sl); add_title(sl, "Literature Survey (2/2)")
add_table(sl, 0.5, 1.3, 12.3, 5.5,
    ["#", "Paper Title", "Authors / Year", "Key Contribution / Relevance"],
    [
        ["6", "Wireless Sensor Networks: A Survey", "Akyildiz et al., Comp. Networks 2002", "Node count scales with area x density (N=AxD). Informs our distributed simulation."],
        ["7", "Edge Computing: Vision and Challenges", "Shi et al., IEEE IoT Journal 2016", "Defines edge computing paradigm. Justifies why encryption must occur at edge, not cloud."],
        ["8", "The Internet of Things: A Survey", "Li et al., Inf. Systems Frontiers 2015", "Comprehensive IoT architecture survey. Validates our sensor data model and edge constraints."],
        ["9", "Report on Lightweight Cryptography", "McKay et al., NIST IR 8114, 2017", "NIST guidelines for lightweight crypto selection. Frames our algorithm comparison criteria."],
        ["10", "Review of Lightweight Block Ciphers", "Hatzivasilis et al., J. Crypto. Eng. 2018", "Systematic review of PRESENT, SPECK, SIMON. Validates our from-scratch implementation approach."],
    ], font_size=11)

# ═══ SLIDE 5: Dataset Generation ═══
sl = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(sl); add_title(sl, "Dataset Generation")
add_bullet_slide(sl, 0.8, 1.5, 11.7, 5.5, [
    "Synthetic IoT Sensor Data:",
    "- We generate realistic JSON payloads mimicking IoT sensor telemetry streams.",
    "- Each record contains: timestamp, device_id, temperature_c, humidity_pct, pressure_hpa, light_lux, sensor_voltage, status.",
    "",
    "Data Sizes Used: 1KB, 10KB, 100KB",
    "- 1KB: Single sensor reading (~5 JSON records) -- typical RFID/smart card payload.",
    "- 10KB: Small batch of readings (~50 records) -- typical for periodic sensor uploads.",
    "- 100KB: Larger aggregated batch (~500 records) -- typical for edge gateway aggregation.",
    "",
    "Generation Strategy:",
    "- Records are accumulated until the serialized JSON meets the target byte size.",
    "- Data is then padded or trimmed to exactly match the requested size.",
    "- The same data is used across all algorithms at each size for fair comparison.",
    "",
    "Total Benchmark Matrix: 4 algorithms x 3 sizes x 5 iterations x 2 ops (enc+dec) = 120 runs"
], size=16)

# ═══ SLIDE 6: Architecture Diagram ═══
sl = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(sl); add_title(sl, "System Architecture")
arch_path = os.path.join(dist_graphs_dir, "architecture_diagram.png")
if os.path.exists(arch_path):
    sl.shapes.add_picture(arch_path, Inches(1.5), Inches(1.3), Inches(10.3), Inches(5.8))
else:
    add_text(sl, 3, 3.5, 7, 1, "[architecture_diagram.png]", size=14, color=GRAY, align=PP_ALIGN.CENTER)

# ═══ SLIDE 7: Execution Flow ═══
sl = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(sl); add_title(sl, "Overall Execution Flow")
add_bullet_slide(sl, 0.8, 1.5, 11.7, 5.5, [
    "Step 1: Data Generation --> generate_iot_data(size_bytes) creates synthetic IoT JSON payloads of exact target size.",
    "",
    "Step 2: Key Generation --> Each algorithm module generates a cryptographic key (AES: 256-bit, ChaCha20: 256-bit, PRESENT: 80-bit, SPECK: 128-bit).",
    "",
    "Step 3: Benchmarked Encryption --> BenchmarkProfiler context manager wraps encrypt(data, key). Captures wall-clock time, CPU%, peak memory, throughput.",
    "",
    "Step 4: Simulated Edge Transmission --> 50ms delay injected between encrypt and decrypt to model real edge-to-cloud network latency.",
    "",
    "Step 5: Benchmarked Decryption --> Same profiler wraps decrypt(nonce, ciphertext, tag, key). Integrity assertion: decrypted == original data.",
    "",
    "Step 6: Results Aggregation --> All 120 measurements saved to CSV/JSON. Pandas aggregation computes averages. Matplotlib generates 6 comparison graphs."
], size=16)

# ═══ SLIDE 8: Algorithm Models ═══
sl = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(sl); add_title(sl, "Encryption Algorithm Models")
add_table(sl, 0.5, 1.3, 12.3, 2.5,
    ["Property", "AES-256-GCM", "ChaCha20-Poly1305", "PRESENT-80", "SPECK-64/128"],
    [
        ["Type", "Block Cipher", "Stream Cipher", "Block Cipher (Lightweight)", "Block Cipher (Lightweight)"],
        ["Structure", "SPN (SubBytes, ShiftRows, MixColumns)", "ARX (Add, Rotate, XOR)", "SPN (4-bit S-Box + P-Box)", "ARX (Add, Rotate, XOR)"],
        ["Block Size", "128 bits", "N/A (stream)", "64 bits", "64 bits"],
        ["Key Size", "256 bits", "256 bits", "80 bits", "128 bits"],
        ["Rounds", "14", "20 quarter-rounds", "31", "27"],
        ["Auth Mode", "GCM (GHASH)", "Poly1305 MAC", "CTR + HMAC-SHA256", "CTR + HMAC-SHA256"],
        ["Gate Count", "~50,000 GE", "~38,000 GE", "~1,570 GE", "~2,400 GE"],
    ], font_size=11)

# ═══ SLIDE 9: Algorithm Methodology (Detail) ═══
sl = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(sl); add_title(sl, "Algorithm Methodology (Detailed)")
add_bullet_slide(sl, 0.8, 1.4, 5.5, 5.5, [
    "AES-256-GCM:",
    "- 14 rounds of SubBytes, ShiftRows, MixColumns, AddRoundKey.",
    "- GCM mode provides authenticated encryption.",
    "- Implemented via PyCryptodome (uses hardware AES-NI if available).",
    "",
    "ChaCha20-Poly1305:",
    "- 20 rounds of quarter-round ARX operations generate a keystream.",
    "- Plaintext XORed with keystream for encryption.",
    "- Poly1305 MAC provides authentication.",
    "- Implemented via PyCryptodome (optimized C backend)."
], size=14)
add_bullet_slide(sl, 6.8, 1.4, 5.5, 5.5, [
    "PRESENT-80 (Built from scratch):",
    "- 31 rounds on 64-bit blocks using 4-bit S-Box substitution + bit permutation (P-Box).",
    "- Wrapped in CTR mode for arbitrary-length payloads.",
    "- HMAC-SHA256 added for integrity (since PRESENT has no built-in authentication).",
    "",
    "SPECK-64/128 (Built from scratch):",
    "- 27 rounds using only bitwise rotation, modular addition, and XOR.",
    "- No S-Boxes: maximizes software efficiency on microcontrollers.",
    "- Wrapped in CTR mode + HMAC-SHA256 for AEAD-like behavior.",
    "- Key expansion generates 27 round keys from the 128-bit master key."
], size=14)

# ═══ SLIDE 10: Results - Throughput & CPU ═══
sl = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(sl); add_title(sl, "Results: Throughput & CPU Usage")
for i, fname in enumerate(["throughput.png", "cpu_usage.png"]):
    path = os.path.join(graphs_dir, fname)
    if os.path.exists(path):
        sl.shapes.add_picture(path, Inches(0.5 + i*6.3), Inches(1.5), Inches(6.0), Inches(3.5))
    else:
        add_text(sl, 0.5 + i*6.3, 2.5, 6.0, 1.0, f"[{fname}]", size=14, color=GRAY, align=PP_ALIGN.CENTER)
add_text(sl, 0.8, 5.3, 11.7, 0.4, "- AES & ChaCha20 achieve ~65-67 MB/s at 100KB. PRESENT and SPECK are orders of magnitude slower in pure Python.", size=15, color=DARK)
add_text(sl, 0.8, 5.8, 11.7, 0.4, "- Lightweight ciphers saturate a single CPU core (~95-100%) due to intensive bit manipulation in Python.", size=15, color=DARK)

# ═══ SLIDE 11: Results - Enc/Dec Time ═══
sl = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(sl); add_title(sl, "Results: Encryption & Decryption Time")
for i, fname in enumerate(["encryption_time.png", "decryption_time.png"]):
    path = os.path.join(graphs_dir, fname)
    if os.path.exists(path):
        sl.shapes.add_picture(path, Inches(0.5 + i*6.3), Inches(1.5), Inches(6.0), Inches(3.5))
    else:
        add_text(sl, 0.5 + i*6.3, 2.5, 6.0, 1.0, f"[{fname}]", size=14, color=GRAY, align=PP_ALIGN.CENTER)
add_text(sl, 0.8, 5.3, 11.7, 0.4, "- PRESENT: ~37s for 100KB encryption. SPECK: ~3s for 100KB. AES/ChaCha20: <2ms for 100KB.", size=15, color=DARK)
add_text(sl, 0.8, 5.8, 11.7, 0.4, "- Time scales linearly with data size for all algorithms, confirming O(n) complexity.", size=15, color=DARK)

# ═══ SLIDE 12: Results - Memory & Dashboard ═══
sl = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(sl); add_title(sl, "Results: Memory & Summary Dashboard")
for i, fname in enumerate(["memory_usage.png", "summary_dashboard.png"]):
    path = os.path.join(graphs_dir, fname)
    if os.path.exists(path):
        sl.shapes.add_picture(path, Inches(0.5 + i*6.3), Inches(1.5), Inches(6.0), Inches(3.5))
    else:
        add_text(sl, 0.5 + i*6.3, 2.5, 6.0, 1.0, f"[{fname}]", size=14, color=GRAY, align=PP_ALIGN.CENTER)
add_text(sl, 0.8, 5.3, 11.7, 0.4, "- At 100KB: AES ~102KB, ChaCha20 ~100KB, PRESENT ~210KB, SPECK ~210KB peak memory.", size=15, color=DARK)

# ═══ SLIDE 13: Results Summary Table (PROPER TABLE) ═══
sl = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(sl); add_title(sl, "Results: Summary Table & Best Algorithm")
add_text(sl, 0.8, 1.3, 11.7, 0.4, "Encryption Performance (averaged across 1KB, 10KB, 100KB -- 5 iterations each):", size=14, bold=True, color=DARK)
add_table(sl, 0.8, 1.8, 11.7, 2.5,
    ["Algorithm", "Avg Time (ms)", "Avg CPU %", "Avg Memory (KB)", "Avg Throughput (MB/s)", "Composite Score"],
    [
        ["AES-256-GCM", "3.811", "65.5", "51.6", "27.14", "7.0"],
        ["ChaCha20-Poly1305", "0.993", "120.0", "39.0", "29.85", "7.0"],
        ["PRESENT-80", "13,783.90", "96.2", "79.1", "0.003", "15.0"],
        ["SPECK-64/128", "1,181.88", "85.7", "79.0", "0.04", "11.0"],
    ], font_size=13)
add_text(sl, 0.8, 4.6, 11.7, 0.4, "Scoring: Each algorithm ranked 1-4 on throughput, latency, CPU, memory. Sum of ranks = composite score. Lower = better.", size=13, color=GRAY)
add_text(sl, 0.8, 5.1, 11.7, 0.5, "BEST OVERALL (tied): AES-256-GCM & ChaCha20-Poly1305 (Score: 7.0)", size=18, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
add_text(sl, 0.8, 5.7, 11.7, 0.4, "- AES wins on CPU (hardware AES-NI). ChaCha20 wins on throughput & memory. On devices without AES-NI, ChaCha20 is the clear winner.", size=14, color=DARK)

# ═══ SLIDE 14: Distributed Scalability ═══
sl = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(sl); add_title(sl, "Results: Scalability & Distributed Performance")
path = os.path.join(dist_graphs_dir, "scalability_dashboard.png")
if os.path.exists(path):
    sl.shapes.add_picture(path, Inches(2.5), Inches(1.5), Inches(8.33), Inches(4.5))
else:
    add_text(sl, 2.5, 3.0, 8.33, 1.0, "[scalability_dashboard.png]", size=14, color=GRAY, align=PP_ALIGN.CENTER)
add_text(sl, 0.8, 6.2, 11.7, 0.5, "System scalability remains nearly linear as the number of nodes increases (N = 5 to 1000).", size=16, color=BLUE, align=PP_ALIGN.CENTER)

# ═══ SLIDE 15: Conclusion - Where Each Algo Excels & Fails ═══
sl = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(sl); add_title(sl, "Conclusion: Where Each Algorithm Excels & Fails")
add_table(sl, 0.5, 1.3, 12.3, 5.5,
    ["Algorithm", "Best For", "Fails When", "Why"],
    [
        ["AES-256-GCM", "High-performance edge nodes\n(Raspberry Pi, x86 with AES-NI)", "Ultra-constrained devices\n(passive RFID, 8-bit MCUs)", "Requires ~50,000 GE and dedicated\nhardware instructions (AES-NI) for speed"],
        ["ChaCha20-Poly1305", "Software-only ARM/RISC-V MCUs\n(no crypto hardware needed)", "Extreme gate-count limits\n(needs ~38,000 GE)", "Pure ARX ops are fast in software but\nstill require moderate silicon area"],
        ["PRESENT-80", "Passive RFID, smart cards\n(needs only ~1,570 GE)", "Any payload >1KB in software\n(~375ms per 1KB, ~37s per 100KB)", "Designed for hardware ASIC, not software.\n31 rounds of 4-bit S-Box are slow in Python."],
        ["SPECK-64/128", "Constrained MCUs needing speed\n(24 cycles/byte on ARM Cortex)", "Large payloads in pure Python\n(~1.7s per 100KB)", "Optimized for C/ASM on MCUs. Python's\ninterpreter overhead negates ARX advantages."],
    ], font_size=12)

# ═══ SLIDE 16: References (1/2) ═══
sl = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(sl); add_title(sl, "References (1/2)")
add_bullet_slide(sl, 0.8, 1.4, 11.7, 5.5, [
    "[1] Daemen, J. & Rijmen, V. (2002). \"The Design of Rijndael: AES.\" Springer.",
    "[2] Bernstein, D.J. (2008). \"ChaCha, a variant of Salsa20.\" SASC 2008.",
    "[3] Bogdanov, A. et al. (2007). \"PRESENT: An Ultra-Lightweight Block Cipher.\" CHES 2007, LNCS 4727.",
    "[4] Beaulieu, R. et al. (2013). \"The SIMON and SPECK Families of Lightweight Block Ciphers.\" IACR ePrint 2013/404.",
    "[5] Thakor, V. et al. (2021). \"Lightweight Cryptography Algorithms for Resource-Constrained IoT Devices.\" IEEE Access, Vol. 9.",
    "[6] Akyildiz, I.F. et al. (2002). \"Wireless Sensor Networks: A Survey.\" Computer Networks, Vol. 38(4).",
    "[7] Shi, W. et al. (2016). \"Edge Computing: Vision and Challenges.\" IEEE Internet of Things Journal, 3(5).",
    "[8] Li, S. et al. (2015). \"The Internet of Things: A Survey.\" Information Systems Frontiers, 17(2).",
    "[9] McKay, K. et al. (2017). \"Report on Lightweight Cryptography.\" NIST Internal Report 8114.",
    "[10] Hatzivasilis, G. et al. (2018). \"Review of Lightweight Block Ciphers.\" Journal of Cryptographic Engineering, 8(2).",
], size=14)

# ═══ SLIDE 16b: References (2/2) ═══
sl = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(sl); add_title(sl, "References (2/2)")
add_bullet_slide(sl, 0.8, 1.4, 11.7, 5.5, [
    "Wikipedia / Online References:",
    "",
    "[W1] https://en.wikipedia.org/wiki/Advanced_Encryption_Standard",
    "[W2] https://en.wikipedia.org/wiki/Salsa20#ChaCha_variant",
    "[W3] https://en.wikipedia.org/wiki/PRESENT_(cipher)"
], size=15)

# ═══ SLIDE 17: Thank You ═══
sl = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(sl)
add_text(sl, 0.8, 2.5, 11.7, 1.5, "Thank You", size=48, bold=True, color=DARK, align=PP_ALIGN.CENTER)
add_line(sl, 4.0)
add_text(sl, 0.8, 4.3, 11.7, 0.5, "Performance Evaluation of Encryption Algorithms on Edge Devices", size=20, color=BLUE, align=PP_ALIGN.CENTER)
add_text(sl, 0.8, 5.0, 11.7, 0.5, "Ready for Q&A", size=24, color=GRAY, align=PP_ALIGN.CENTER)

# Save
out = os.path.join(os.path.dirname(__file__), "Encryption_Edge_Devices_Presentation.pptx")
prs.save(out)
print(f"PPT saved to: {out}")
